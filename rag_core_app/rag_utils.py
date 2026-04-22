import sys
if sys.stdout.encoding.lower() != 'utf-8':
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except AttributeError:
        pass

import os
import shutil
import pytesseract
import requests
import concurrent.futures
from io import BytesIO
from datetime import datetime
from pathlib import Path

BASE_DIR = Path(__file__).resolve().parent.parent

from PIL import Image
from pdf2image import convert_from_path
from langchain_community.document_loaders import (
    PyMuPDFLoader,
    TextLoader,
    Docx2txtLoader,
    CSVLoader,
    WebBaseLoader
)
from langchain.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_groq import ChatGroq
from .models import ChatMessage
from langchain_community.tools import DuckDuckGoSearchResults
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser

EMBEDDING_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"

try:
    GLOBAL_EMBEDDINGS = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL_NAME)
    CHAT_LLM = ChatGroq(temperature=0.2, model_name="llama-3.1-8b-instant", streaming=True)
    ROUTER_LLM = ChatGroq(temperature=0.0, model_name="llama-3.1-8b-instant")
except Exception:
    GLOBAL_EMBEDDINGS = None
    CHAT_LLM = None
    ROUTER_LLM = None

HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
}


def get_db_path(session_id):
    """Returns the absolute path to the session's FAISS index directory."""
    if not session_id:
        return None
    return str(BASE_DIR / "faiss_indexes" / f"session_{session_id}")


def load_single_file(file_path):
    """Loads and extracts text from a URL or local file of any supported type."""
    try:
        if file_path.startswith("http://") or file_path.startswith("https://"):
            if any(file_path.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.webp']):
                try:
                    response = requests.get(file_path, headers=HEADERS, timeout=10)
                    text = pytesseract.image_to_string(Image.open(BytesIO(response.content)))
                    if text.strip():
                        return [Document(page_content=text, metadata={"source": file_path, "type": "image_url"})]
                    return []
                except Exception:
                    return []

            try:
                return WebBaseLoader(file_path, header_template=HEADERS).load()
            except Exception:
                return []

        ext = os.path.splitext(file_path)[1].lower()

        if ext in ['.png', '.jpg', '.jpeg']:
            try:
                text = pytesseract.image_to_string(Image.open(file_path))
                if text.strip():
                    return [Document(page_content=text, metadata={"source": file_path})]
                return []
            except Exception:
                return []

        if ext == ".pdf":
            try:
                return PyMuPDFLoader(file_path).load()
            except Exception as e:
                print(f"[WARN] PDF load error ({file_path}): {e}")
                return []

        elif ext == ".txt":
            # Try common encodings in order; Windows Notepad defaults to UTF-16.
            for enc in ("utf-8-sig", "utf-16", "utf-16-le", "latin-1"):
                try:
                    with open(file_path, "r", encoding=enc) as fh:
                        text = fh.read()
                    if text.strip():
                        return [Document(page_content=text, metadata={"source": file_path})]
                    break
                except (UnicodeDecodeError, UnicodeError):
                    continue
                except Exception as e:
                    print(f"[WARN] TXT load error ({file_path}): {e}")
                    break
            return []

        elif ext in [".docx", ".doc"]:
            try:
                return Docx2txtLoader(file_path).load()
            except Exception as e:
                print(f"[WARN] DOCX load error ({file_path}): {e}")
                return []

        elif ext == ".csv":
            try:
                return CSVLoader(file_path, encoding='utf-8').load()
            except Exception:
                try:
                    return CSVLoader(file_path, encoding='latin-1').load()
                except Exception as e:
                    print(f"[WARN] CSV load error ({file_path}): {e}")
                    return []

        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                slides_text = []
                for i, slide in enumerate(prs.slides):
                    slide_content = [
                        shape.text.strip()
                        for shape in slide.shapes
                        if hasattr(shape, "text") and shape.text.strip()
                    ]
                    if slide_content:
                        slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_content))
                full_text = "\n\n".join(slides_text)
                if full_text.strip():
                    return [Document(page_content=full_text, metadata={"source": file_path, "type": "pptx"})]
                return []
            except Exception as e:
                print(f"[WARN] PPTX load error ({file_path}): {e}")
                return []

        elif ext in [".xlsx", ".xls"]:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                all_rows = []
                for ws in wb.worksheets:
                    all_rows.append(f"--- Sheet: {ws.title} ---")
                    for row in ws.iter_rows(values_only=True):
                        cleaned = [str(cell) for cell in row if cell is not None]
                        if cleaned:
                            all_rows.append(" | ".join(cleaned))
                full_text = "\n".join(all_rows)
                if full_text.strip():
                    return [Document(page_content=full_text, metadata={"source": file_path, "type": "xlsx"})]
                return []
            except Exception as e:
                print(f"[WARN] XLSX load error ({file_path}): {e}")
                return []

        else:
            try:
                return TextLoader(file_path, autodetect_encoding=True).load()
            except Exception as e:
                print(f"[WARN] Unknown file type load error ({file_path}): {e}")
                return []

    except Exception:
        return []


def process_files_bulk(file_paths, session_id):
    """Ingests a list of file paths and/or URLs into the session's FAISS index."""
    if not session_id or not GLOBAL_EMBEDDINGS or not file_paths:
        return False

    db_path = get_db_path(session_id)
    all_documents = []

    with concurrent.futures.ThreadPoolExecutor() as executor:
        for res in executor.map(load_single_file, file_paths):
            if res:
                all_documents.extend(res)

    if not all_documents:
        return False

    docs = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100).split_documents(all_documents)
    if not docs:
        return False

    try:
        if os.path.exists(db_path):
            try:
                vector_store = FAISS.load_local(db_path, GLOBAL_EMBEDDINGS, allow_dangerous_deserialization=True)
                vector_store.add_documents(docs)
                vector_store.save_local(db_path)
            except Exception as e:
                print(f"[WARN] FAISS merge failed, rebuilding index: {e}")
                FAISS.from_documents(docs, GLOBAL_EMBEDDINGS).save_local(db_path)
        else:
            FAISS.from_documents(docs, GLOBAL_EMBEDDINGS).save_local(db_path)
        return True
    except Exception as e:
        print(f"[ERROR] FAISS indexing failed: {e}")
        return False


def perform_web_search(query):
    try:
        return DuckDuckGoSearchResults(num_results=4).run(query)
    except Exception as e:
        return f"Web search failed: {e}"


def get_answer(query, session_id):
    db_path = get_db_path(session_id)
    current_date = datetime.now().strftime("%Y-%m-%d")

    recent_history = ChatMessage.objects.filter(session_id=session_id).order_by('-timestamp')[:6]
    history_text = "\n".join(
        [f"{'User' if msg.is_user else 'AI'}: {msg.text}" for msg in reversed(recent_history)]
    )

    try:
        router_prompt = ChatPromptTemplate.from_template(
            "Classify the following user input: '{question}'.\n"
            "If the user is asking for ANY real-world facts, current events, knowledge, programming help, or specific information, reply ONLY with 'QUERY'.\n"
            "If the user is ONLY making small talk, greeting you, or saying thanks, reply ONLY with 'CHAT'."
        )
        intent = (router_prompt | ROUTER_LLM | StrOutputParser()).invoke({"question": query}).strip().upper()
    except Exception:
        intent = "QUERY"

    if "CHAT" in intent:
        prompt = ChatPromptTemplate.from_template(
            """You are a highly capable, precise, and professional AI assistant.
            Current Date: {date}

            Conversation History:
            {history}

            User: {question}

            Reply naturally, concisely, and use clean markdown formatting. Do not use filler language."""
        )
        for chunk in (prompt | CHAT_LLM | StrOutputParser()).stream({"date": current_date, "history": history_text, "question": query}):
            yield chunk
        return

    has_index = db_path and os.path.exists(db_path) and GLOBAL_EMBEDDINGS
    if has_index:
        try:
            vector_store = FAISS.load_local(db_path, GLOBAL_EMBEDDINGS, allow_dangerous_deserialization=True)
            results = vector_store.similarity_search_with_score(query, k=5)
            if results:
                context_text = "\n\n".join([d.page_content for d, _ in results])
                source_type = "Uploaded Document"
                print(f"[RAG] Session {session_id}: {len(results)} chunks retrieved. Top score: {results[0][1]:.4f}")
            else:
                context_text = perform_web_search(query)
                source_type = "Web Search"
        except Exception as e:
            print(f"[WARN] Retrieval error for session {session_id}: {e}")
            context_text = perform_web_search(query)
            source_type = "Web Search"
    else:
        context_text = perform_web_search(query)
        source_type = "Web Search"

    system_prompt = """You are an expert, precision-focused AI assistant.
    Current Date: {date}. Always use this date context to ensure your answers are up-to-date and relevant.

    Conversation History:
    {history}

    Context ({source}):
    {context}

    User Question: {question}

    Instructions:
    1. Precision & Clarity: Provide direct, concise, and highly accurate answers. Eliminate fluff, repetitive intros, and filler words.
    2. Clean Formatting: Heavily utilize Markdown. Structure your answers with clear headings (`###`), bullet points, and **bold text** for key terms.
    3. Knowledge Strategy: Base your answers STRICTLY on the provided Context. If the context does not contain the answer, use your General Knowledge only if you are absolutely certain. If uncertain, state: "I don't know. Please upload relevant documents for this topic."
    4. Sources: At the very end of your response, cite sources as inline Markdown links: `[Source 1](URL1) [Source 2](URL2)`. Do not use a bulleted list."""

    for chunk in (ChatPromptTemplate.from_template(system_prompt) | CHAT_LLM | StrOutputParser()).stream({
        "date": current_date,
        "source": source_type,
        "history": history_text,
        "context": context_text,
        "question": query
    }):
        yield chunk


def clear_data(session_id):
    db_path = get_db_path(session_id)
    if db_path and os.path.exists(db_path):
        try:
            shutil.rmtree(db_path)
        except Exception:
            pass


def generate_chat_title(user_message, bot_response):
    try:
        llm = ChatGroq(temperature=0.3, model_name="llama-3.1-8b-instant")
        prompt = (
            "Determine the core topic of this conversation and create a concise 2-4 word title for it. "
            "Return ONLY the title with no quotes or extra text.\n"
            f"User: {user_message}\nAI: {bot_response}"
        )
        return llm.invoke(prompt).content.strip().replace('"', '')[:50]
    except Exception:
        return user_message[:30]